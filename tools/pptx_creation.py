#!/usr/bin/env python3
"""
pptx_creation.py

Iterate through scenes & configs, pick up the pre-stitched views for each,
and dump them into a PowerPoint or PDF (one slide per scene+config).
"""
import os
import json
import sys
sys.path.insert(0, os.getcwd())  # Ensure the current directory is in the path

from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib as mpl
from PIL import Image
from arguments import get_args
from utils import get_scenes_list, get_out_dir, log_paths, read_json, make_scene_legend
from gpt_inference import parse_material_list
from plot.plot_utils import make_legend
import open_clip
from feature_fusion import CLIP_BACKBONE, CLIP_CHECKPOINT
from predict_property import predict_physical_property_integral


from PIL import Image, ImageDraw, ImageFont
import textwrap
import os

def estimate_max_font_size(target_h, n_rows, pad=20, line_gap=10, min_font=18, max_font=40):
    """
    Estimate the largest font size such that all rows fit in the given target height.
    """
    for font_size in range(max_font, min_font - 1, -1):
        total_h = n_rows * (font_size + line_gap) + 2 * pad
        if total_h <= target_h:
            return font_size
    return min_font

# ------------------------------------------------------------------
# Utility: build a tiny “spec sheet” image from the three lists
# ------------------------------------------------------------------
from PIL import Image, ImageDraw, ImageFont

def build_material_note(mat_names, mat_vals, mat_tns,
                        target_h,
                        max_font=90,      # upper bound (pt)
                        min_font=18,      # lower bound (pt)
                        pad=20, line_gap=10,
                        title='Material        ρ [kg/m³]       t [cm]',
                        gt_mass = None, pred_mass = None, args=None, N = None, caption=None, heavy_str=None):
    """
    Create a note image whose font size is chosen to *fill* the target
    height.  You just pass in the height of your legend or main panel.
    """
    n_rows = len(mat_names) + 1                       # +1 for the header

    # Pick the largest font that still fits vertically
    best_font_size = max_font
    while best_font_size >= min_font:
        total_h = (best_font_size + line_gap) * n_rows + 2 * pad
        if total_h <= target_h:
            break
        best_font_size -= 2                    # step down
    else:
        best_font_size = min_font              # fallback

    # Load a font
    try:
        font = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            best_font_size
        )
    except IOError:
        font = ImageFont.load_default()

    # Compose row strings (“min–max” for ranges)
    rows = [title]
    for n, v_range, t_range in zip(mat_names, mat_vals, mat_tns):
        v_lo, v_hi = v_range
        t_lo, t_hi = t_range
        rows.append(f"{n:<16} {v_lo:4.0f}–{v_hi:<4.0f}    {t_lo:.2f}–{t_hi:.2f}")

    rows.append(f'Correction Factor: {args.correction_factor}')
    rows.append(f'feature load name: {args.feature_load_name}')
    rows.append(f'Number of dense points: {N}')

    if caption:
        # trim caption if too long
        if len(caption) > 100:
            caption = caption[:97] + '...'
        rows.append(f'Caption: {caption}')
    
    if heavy_str:
        
        rows.append(f'Heavy state: {heavy_str}')
    # Width just big enough for the longest row
    max_w = max(font.getbbox(txt)[2] for txt in rows) + 2 * pad
    note = Image.new("RGB", (max_w, target_h), (255, 255, 255))
    draw = ImageDraw.Draw(note)

    # Draw text
    y = pad
    for txt in rows:
        draw.text((pad, y), txt, fill=(0, 0, 0), font=font)
        y += best_font_size + line_gap
    # write down gt and pred mass as well (pred_mass is a list of 2 ele)

    if gt_mass is not None and pred_mass is not None:
        gt_txt = f"GT Mass: {gt_mass:.4f} kg"
        pred_txt = f"Pred Mass: {pred_mass[0]:.4f}–{pred_mass[1]:.4f} kg"
        draw.text((pad, y), gt_txt, fill=(0, 0, 0), font=font)
        y += best_font_size + line_gap
        draw.text((pad, y), pred_txt, fill=(0, 0, 0), font=font)
    return note


def main():
    args = get_args()

    # Mode: pptx (default) or pdf
    export_mode = getattr(args, "slide_mode", "pptx").lower()
    if export_mode not in ("pptx", "pdf"):
        raise ValueError(f"Invalid mode: {export_mode}. Choose 'pptx' or 'pdf'.")

    # 1) Get scenes
    scenes = get_scenes_list(args)

    # 2) Set up
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_basename = os.path.join(args.pptx_path, f"{timestamp}")

    if export_mode == "pptx":
        prs = Presentation()
        blank = prs.slide_layouts[6]
    else:
        pdf_images = []

    print(f"🧪 Building {export_mode.upper()} for {len(scenes)} scenes…")

    gt_mass_json = read_json(args.gt_json) if args.gt_json else None

    clip_model, _, preprocess = open_clip.create_model_and_transforms(CLIP_BACKBONE, pretrained=CLIP_CHECKPOINT)
    clip_model.to(args.device)
    clip_tokenizer = open_clip.get_tokenizer(CLIP_BACKBONE)

    # 3) Process each scene
    for scene in scenes:
        args.scene_name = scene
        scene_dir = os.path.join(args.data_dir, 'scenes', scene)
        old_feature_load_name = args.feature_load_name
        if old_feature_load_name == '2d_patch':
            args.feature_load_name = 'default'
            args.feature_save_name = 'default'
        out_dir = get_out_dir(args)

        if old_feature_load_name == '2d_patch':
            args.feature_load_name = old_feature_load_name
            args.feature_save_name = old_feature_load_name
        
        scene_gt_name = scene.split('_')[0]

        gt_mass = gt_mass_json[scene_gt_name] if gt_mass_json else None

        total_val, debug = predict_physical_property_integral(
            args, scene_dir, clip_model, clip_tokenizer,
        )

        #stitched_path = os.path.join(out_dir + "_cuboid", f"{args.viz_save_name}_stitched_all_views.png")
        stitched_path = os.path.join(out_dir + f"_cuboid_{args.plot_cuboids_mode}", f"{args.viz_save_name}_stitched_all_views.png")
        legend_path = os.path.join(out_dir, '_legend.png')
        combined_path = os.path.join(out_dir + f'_{args.feature_load_name}', '_combined_with_legend.png')

        os.makedirs(os.path.dirname(combined_path), exist_ok=True)

        if not os.path.exists(stitched_path) and args.plot_type != 'skip_stitched':
            print(f"[❌] Missing stitched image: {stitched_path}")
            continue

        mat_names, mat_vals, mat_tns, caption, heavy_str = make_scene_legend(args, scene_dir, out_dir)

        if args.plot_type != 'skip_stitched':
            main_img = Image.open(stitched_path).convert("RGB")
        else:
            #scene_dir = os.path.join(args.data_dir, 'scenes', scene)
            #print('path:', os.path.join(scene_dir, 'images', f"{scene}_00.png"))
            main_img = Image.open(os.path.join(scene_dir, 'images', f"{scene}_00.png")).convert("RGB")

        legend_img = Image.open(legend_path).convert("RGB")
        h_target = max(main_img.height, legend_img.height)

        #print('legend path:', legend_path)
        #print('legend img shape:', legend_img.size)
        n_materials = len(mat_names)
        n_rows = n_materials + 1  # 1 title row
        if gt_mass is not None and total_val is not None:
            n_rows += 2  # GT + pred mass
            
        #max_font_size = estimate_max_font_size(h_target, n_rows)

        #print('max_font_size:', max_font_size)
        note_img = build_material_note(mat_names, mat_vals, mat_tns, target_h=h_target, gt_mass=gt_mass, pred_mass=total_val, args=args, N = len(debug['dense_pts']), caption=caption, heavy_str=heavy_str)

        #print('ehehehhe')
        def pad_to_height(im, target_h):
            if im.height == target_h:
                return im
            new_im = Image.new("RGB", (im.width, target_h), (255, 255, 255))
            new_im.paste(im, (0, 0))
            return new_im

        main_img   = pad_to_height(main_img,   h_target)
        legend_img = pad_to_height(legend_img, h_target)
        note_img   = pad_to_height(note_img,   h_target)


        #print('hehehe')
        # # Pad heights to match
        # h_main, h_legend = main_img.height, legend_img.height
        # if h_main > h_legend:
        #     new_legend = Image.new("RGB", (legend_img.width, h_main), (255, 255, 255))
        #     new_legend.paste(legend_img, (0, 0))
        #     legend_img = new_legend
        # elif h_legend > h_main:
        #     new_main = Image.new("RGB", (main_img.width, h_legend), (255, 255, 255))
        #     new_main.paste(main_img, (0, 0))
        #     main_img = new_main
        
        #Eif args.plot_type != 'skip_stitched':
        images = [main_img, legend_img, note_img]
        
        if not os.path.exists(combined_path) or args.overwrite:
            # Combine horizontally
            # combined_img = Image.new("RGB", (main_img.width + legend_img.width, main_img.height), (255, 255, 255))
            # combined_img.paste(main_img, (0, 0))
            # combined_img.paste(legend_img, (main_img.width, 0))
            # combined_img.save(combined_path)
            #print(f"🖼️ Saved stitched image with legend to: {combined_path}")
            # combined_w = main_img.width + legend_img.width + note_img.width
            # combined_img = Image.new("RGB", (combined_w, h_target), (255, 255, 255))
            # x = 0
            # for im in (main_img, legend_img, note_img):
            #     combined_img.paste(im, (x, 0))
            #     x += im.width
            combined_w = 0
            for im in images:
                combined_w += im.width

            combined_img = Image.new("RGB", (combined_w, h_target), (255, 255, 255))

            x = 0
            for im in images:
                combined_img.paste(im, (x, 0))
                x += im.width
            combined_img.save(combined_path)
            print(f"🖼️ Saved stitched image with legend + spec sheet to: {combined_path}")
        else:
            combined_img = Image.open(combined_path).convert("RGB")
            combined_img.close()
            print(f"🖼️ Using existing combined image: {combined_path}")
        
        # combined_w = main_img.width + legend_img.width + note_img.width
        # combined_img = Image.new("RGB", (combined_w, h_target), (255, 255, 255))
        # x = 0
        # for im in (main_img, legend_img, note_img):
        #     combined_img.paste(im, (x, 0))
        #     x += im.width


        if export_mode == "pptx":
            slide = prs.slides.add_slide(blank)
            # Add title
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            para = txbox.text_frame.paragraphs[0]
            para.text = f"{scene} — {args.feature_load_name}"
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            slide.shapes.add_picture(combined_path, Inches(0.5), Inches(1), width=Inches(10.5))
        else:
            pdf_images.append(combined_img)


    # 4) Save output
    if export_mode == "pptx":
        out_file = output_basename + ".pptx"
        prs.save(out_file)
    else:
        out_file = output_basename + ".pdf"
        if pdf_images:
            pdf_images[0].save(out_file, save_all=True, append_images=pdf_images[1:])
        else:
            print("[❌] No images to save as PDF.")
            return

    print(f"[📊] Saved to {out_file}")
    log_paths(args.log_file, [out_file])

if __name__ == "__main__":
    main()
