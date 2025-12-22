import sys, os 
sys.path.insert(0, os.getcwd())


import os, json
import numpy as np
from scipy.spatial import cKDTree
import matplotlib.pyplot as plt
import csv

from utils import *
from arguments import get_args

from external.PhysX_3D.plot_utils import save_point_cloud_plot
from tqdm import tqdm

from external.splatter_image.utils.general_utils import pca_pose, size_from_pose_centers
from external.PhysX_3D.plot_utils import plot_series

from pptx import Presentation
from pptx.util import Inches

from datetime import datetime

LIST_90 = [0, 1, 6, 11, 20, 37, 42, 45, 51, 67, 70, 89, 90]
LIST_60 = [33, 66]
STAY = [10, 15, 20, 23, 25, 27, 33, 45, 48, 52, 54, 61, 62, 64, 76, 92]
def _nn_dist_idx(src: np.ndarray, dst: np.ndarray):
    """
    For each point in `src`, find its nearest neighbor in `dst`.
    Returns (distances, indices) where distances are Euclidean.
    """
    if src.size == 0 or dst.size == 0:
        raise ValueError("Empty point cloud passed to nearest-neighbor search.")
    tree = cKDTree(dst)
    dists, idx = tree.query(src, k=1, workers=-1)
    return dists, idx


def chamfer_distance(A: np.ndarray, B: np.ndarray, squared: bool = False) -> float:
    """
    Chamfer Distance between point sets A and B.
    CD(A,B) = (1/|A|) * sum_{a in A} min_{b in B} ||a-b||  + (1/|B|) * sum_{b in B} min_{a in A} ||b-a||
    If `squared=True`, use squared distances (be explicit and consistent across comparisons).
    """
    da, _ = _nn_dist_idx(A, B)   # each A -> nearest in B
    db, _ = _nn_dist_idx(B, A)   # each B -> nearest in A
    if squared:
        da = da**2
        db = db**2
    return da.mean() + db.mean()


def fscore_at_tau(A: np.ndarray, B: np.ndarray, tau: float):
    """
    F-score@τ comparing coverage within a distance threshold tau.
      Precision = fraction of A within τ of B
      Recall    = fraction of B within τ of A
      F         = 2PR / (P+R)  (0 if both P and R are 0)
    Returns (precision, recall, fscore).
    """
    if tau <= 0:
        raise ValueError("tau must be positive.")

    da, _ = _nn_dist_idx(A, B)  # distances A -> B
    db, _ = _nn_dist_idx(B, A)  # distances B -> A

    precision = float((da <= tau).mean()) if len(da) else 0.0
    recall    = float((db <= tau).mean()) if len(db) else 0.0
    fscore    = (2 * precision * recall / (precision + recall)) if (precision + recall) > 0 else 0.0
    return precision, recall, fscore

def translate_to_origin(pts):
    """Translate point cloud to have centroid at the origin."""
    centroid = np.mean(pts, axis=0)
    return pts - centroid

def normal_consistency(
    A: np.ndarray, NA: np.ndarray,
    B: np.ndarray, NB: np.ndarray,
    eps: float = 1e-9
) -> float:
    """
    Normal Consistency (NC) between two oriented point sets.
      NC = 0.5 * ( mean_a |n_a · n_{b(a)}| + mean_b |n_b · n_{a(b)}| )
    where b(a) is A's nearest neighbor in B, and a(b) vice versa.
    Normals are normalized internally. Returns a value in [0, 1] where higher is better.
    """
    if len(A) != len(NA) or len(B) != len(NB):
        raise ValueError("Points and normals must have matching lengths for each set.")

    # Normalize normals safely
    def _normalize(N):
        nrm = np.linalg.norm(N, axis=1, keepdims=True)
        nrm = np.maximum(nrm, eps)
        return N / nrm

    NAu = _normalize(NA)
    NBu = _normalize(NB)

    # A -> B
    _, idx_ab = _nn_dist_idx(A, B)
    dots_ab = np.abs((NAu * NBu[idx_ab]).sum(axis=1))

    # B -> A
    _, idx_ba = _nn_dist_idx(B, A)
    dots_ba = np.abs((NBu * NAu[idx_ba]).sum(axis=1))

    return 0.5 * (dots_ab.mean() + dots_ba.mean())


# --- your metric fns from above here (_nn_dist_idx, chamfer_distance, fscore_at_tau, normal_consistency) ---

def compute_scene_metrics(pts_a, pts_b, normals_a=None, normals_b=None, tau=0.1):
    """Return a dict of metrics for one scene."""
    cd = chamfer_distance(pts_a, pts_b, squared=False)
    p, r, f = fscore_at_tau(pts_a, pts_b, tau=tau)

    nc = None
    if normals_a is not None and normals_b is not None:
        nc = normal_consistency(pts_a, normals_a, pts_b, normals_b)

    return {"chamfer": float(cd), "precision": float(p), "recall": float(r), "fscore": float(f), "normal_consistency": (None if nc is None else float(nc))}

def save_metrics_csv(per_scene_metrics, save_path):
    """Write a tidy CSV for later analysis."""
    keys = ["scene", "chamfer", "precision", "recall", "fscore", "normal_consistency"]
    with open(save_path, "w", newline="") as f:
        w = csv.DictWriter(f, fieldnames=keys)
        w.writeheader()
        for row in per_scene_metrics:
            w.writerow({k: row.get(k, None) for k in keys})

def _safe_normals_load(path):
    return np.loadtxt(path) if os.path.exists(path) else None


def _subsample(points: np.ndarray, k: int) -> np.ndarray:
    """Stride-based downsample: keep every k-th point (k>=1)."""
    if k is None or k <= 1:
        return points
    return points[::k, :]

def _best_rigid_transform(A: np.ndarray, B: np.ndarray):
    """
    Compute best-fit rigid transform (R, t) minimizing || R*B + t - A ||_F
    using SVD (A,B are Nx3 and correspond one-to-one).
    Returns R (3x3), t (3,), and RMS error.
    """
    assert A.shape == B.shape and A.shape[1] == 3
    mu_A = A.mean(axis=0)
    mu_B = B.mean(axis=0)
    AA = A - mu_A
    BB = B - mu_B
    H = BB.T @ AA
    U, S, Vt = np.linalg.svd(H)
    R = Vt.T @ U.T
    # Handle reflection
    if np.linalg.det(R) < 0:
        Vt[-1, :] *= -1
        R = Vt.T @ U.T
    t = mu_A - R @ mu_B
    rms = np.sqrt(np.mean(np.sum((A - (B @ R.T + t))**2, axis=1)))
    return R, t, rms

def icp_align_to_source(
    source: np.ndarray,
    others: list,
    downsample_factor: int = 1,
    max_iters: int = 50,
    tol: float = 1e-6,
    trim_fraction: float = 0.0,
):
    """
    Align each point cloud in `others` to `source` using point-to-point ICP.

    Args:
        source: (Ns,3) float array.
        others: list of float arrays [(N1,3), (N2,3), ...].
        downsample_factor: if >1, stride-subsample both source and target by this factor.
        max_iters: max ICP iterations.
        tol: stop if change in mean squared error improves by < tol.
        trim_fraction: 0.0..0.5. If >0, drop largest residuals each iter (robustness).

    Returns:
        aligned_list: list of aligned full-res clouds (same shapes as input).
        transforms:   list of 4x4 homogeneous transforms mapping original -> aligned
        stats:        list of dicts with per-target 'iters', 'rms', 'converged'
    """
    # Downsample once for the source used in correspondence
    src_ds = _subsample(np.asarray(source, dtype=np.float64), downsample_factor)
    tree = cKDTree(src_ds)

    aligned_list, transforms, stats = [], [], []

    for tgt in others:
        tgt_full = np.asarray(tgt, dtype=np.float64)
        tgt_ds = _subsample(tgt_full, downsample_factor)

        # Running transform (R,t) accumulates over iterations
        R_acc = np.eye(3)
        t_acc = np.zeros(3)

        prev_mse = np.inf
        converged = False
        it_used = 0

        X = tgt_ds.copy()  # working copy (transformed target, downsampled)

        for it in range(max_iters):
            # 1) Correspondences: nearest neighbors in source
            dists, idx = tree.query(X, k=1, workers=-1)
            Y = src_ds[idx]

            # Optional trimming to reduce outlier influence
            if trim_fraction > 0.0:
                m = len(dists)
                keep = int(np.ceil((1.0 - trim_fraction) * m))
                order = np.argpartition(dists, keep-1)[:keep]
                X_m = X[order]
                Y_m = Y[order]
            else:
                X_m, Y_m = X, Y

            # 2) Best-fit rigid transform
            R, t, rms = _best_rigid_transform(Y_m, X_m)  # maps X_m -> Y_m
            # 3) Update working cloud and accumulate transform
            X = (X @ R.T) + t

            # Compose transforms: new overall transform maps original -> current
            R_acc = R @ R_acc
            t_acc = R @ t_acc + t

            mse = np.mean(dists**2)
            if abs(prev_mse - mse) < tol:
                converged = True
                it_used = it + 1
                break
            prev_mse = mse
            it_used = it + 1

        # Apply final transform to the FULL-RES target
        aligned = (tgt_full @ R_acc.T) + t_acc

        # Homogeneous 4x4
        T = np.eye(4)
        T[:3, :3] = R_acc
        T[:3, 3] = t_acc

        aligned_list.append(aligned)
        transforms.append(T)
        stats.append({"iters": it_used, "rms": float(np.sqrt(prev_mse)), "converged": converged})

    return aligned_list, transforms, stats

from scipy.spatial import cKDTree

def _pca_frame(points: np.ndarray, up_ref=np.array([0.0, 1.0, 0.0])):
    """
    Compute an orthonormal, right-handed frame from PCA where:
      - y-axis = smallest-variance axis (vertical for many household objects)
      - x-axis = largest-variance axis
      - z-axis = x × y  (then x := y × z to enforce perfect orthonormality)
      - Flip y to align with +up_ref if needed.
    Returns (R, centroid) with R as 3x3 rotation whose columns are [x, y, z].
    """
    P = np.asarray(points, dtype=np.float64)
    c = P.mean(axis=0)
    X = P - c

    # Covariance & eigendecomposition
    C = (X.T @ X) / max(len(X) - 1, 1)
    evals, evecs = np.linalg.eigh(C)   # eigh: symmetric -> sorted ascending
    # Sort by variance descending: a1 (largest), a2, a3 (smallest)
    order = np.argsort(evals)[::-1]
    a1, a2, a3 = evecs[:, order].T  # rows as vectors

    # Choose axes: y = smallest variance (a3), x = largest (a1)
    y = a3 / np.linalg.norm(a3)
    # Make y point along +up_ref if provided
    if np.dot(y, up_ref) < 0:
        y = -y

    x = a1 - np.dot(a1, y) * y  # remove any tiny component along y
    x /= np.linalg.norm(x) + 1e-12

    z = np.cross(x, y)
    z /= np.linalg.norm(z) + 1e-12

    # Re-orthogonalize x to ensure perfect right-handedness
    x = np.cross(y, z)
    x /= np.linalg.norm(x) + 1e-12

    R = np.stack([x, y, z], axis=1)  # columns
    # Ensure right-handed (det > 0); if not, flip z
    if np.linalg.det(R) < 0:
        z = -z
        R = np.stack([x, y, z], axis=1)

    return R, c


def _icp_once(source, target, max_iters=40, tol=1e-6, trim_fraction=0.0):
    """
    Point-to-point ICP aligning target -> source. Returns (R,t).
    """
    src = np.asarray(source, dtype=np.float64)
    tgt = np.asarray(target, dtype=np.float64)

    tree = cKDTree(src)
    R_acc = np.eye(3)
    t_acc = np.zeros(3)
    prev_mse = np.inf

    X = tgt.copy()
    for _ in range(max_iters):
        dists, idx = tree.query(X, k=1, workers=-1)
        Y = src[idx]

        if trim_fraction > 0:
            m = len(dists)
            keep = int(np.ceil((1.0 - trim_fraction) * m))
            sel = np.argpartition(dists, keep-1)[:keep]
            X_m, Y_m = X[sel], Y[sel]
        else:
            X_m, Y_m = X, Y

        R, t, _ = _best_rigid_transform(Y_m, X_m)
        X = (X @ R.T) + t

        R_acc = R @ R_acc
        t_acc = R @ t_acc + t

        mse = float((dists**2).mean())
        if abs(prev_mse - mse) < tol:
            break
        prev_mse = mse

    return R_acc, t_acc

def intensive_alignment(
    source: np.ndarray,
    other: np.ndarray,
    downsample_factor: int = 10,
    do_icp: bool = True,
    icp_trim: float = 0.2,
    icp_max_iters: int = 40,
):
    """
    Coarse PCA alignment (respecting +Y up & right-handedness) + optional ICP refinement.
    Aligns `other` to `source`.

    Args:
        source: (Ns,3)
        other:  (Nt,3)
        downsample_factor: >1 to speed ICP (keeps every k-th point)
        do_icp: run ICP refinement after PCA coarse alignment
        icp_trim: 0..0.5 - drop top fraction of largest residuals per iteration
        icp_max_iters: ICP iterations

    Returns:
        aligned_full: (Nt,3) other aligned into source frame
        T_total: 4x4 transform mapping original `other` -> aligned
        T_pca:   4x4 coarse transform from PCA frames
        T_icp:   4x4 refinement (identity if do_icp=False)
    """
    S = np.asarray(source, dtype=np.float64)
    O = np.asarray(other, dtype=np.float64)

    # ----- PCA frames (coarse) -----
    R_s, c_s = _pca_frame(S)     # frame of source
    R_o, c_o = _pca_frame(O)     # frame of other

    # Rotation that maps other's frame to source's frame
    R_coarse = R_s @ R_o.T
    t_coarse = c_s - R_coarse @ c_o

    # Apply coarse alignment to full-res "other"
    O_coarse = (O @ R_coarse.T) + t_coarse

    # Prepare 4x4 for coarse
    T_pca = np.eye(4)
    T_pca[:3, :3] = R_coarse
    T_pca[:3, 3] = t_coarse

    # ----- Optional ICP refinement -----
    T_icp = np.eye(4)
    if do_icp:
        # Downsample for speed
        S_ds = S[::max(1, downsample_factor)]
        O_ds = O_coarse[::max(1, downsample_factor)]

        R_ref, t_ref = _icp_once(S_ds, O_ds, max_iters=icp_max_iters, tol=1e-6, trim_fraction=icp_trim)

        # Compose total transform: T_total = T_ref ∘ T_pca
        O_aligned = (O_coarse @ R_ref.T) + t_ref

        T_icp[:3, :3] = R_ref
        T_icp[:3, 3] = t_ref
    else:
        O_aligned = O_coarse

    # Compose total
    T_total = np.eye(4)
    T_total[:3, :3] = T_icp[:3, :3] @ T_pca[:3, :3]
    T_total[:3, 3]  = T_icp[:3, :3] @ T_pca[:3, 3] + T_icp[:3, 3]

    return O_aligned, T_total, T_pca, T_icp

def _aabb_diagonal_length(pts: np.ndarray) -> float:
    if pts.size == 0:
        return 0.0
    # Ignore NaNs/Infs if any slipped in
    pts = pts[np.isfinite(pts).all(axis=1)]
    if pts.size == 0:
        return 0.0
    mn = pts.min(axis=0)
    mx = pts.max(axis=0)
    return float(np.linalg.norm(mx - mn, ord=2))

def translate_trimmed_mean(source_pts, pred_pts, trim_percentage=0.1):
    """
    Use trimmed mean - remove top/bottom percentiles before computing mean.
    """
    def trimmed_mean(points, trim_pct):
        n_trim = int(len(points) * trim_pct)
        center = np.zeros(3)
        for dim in range(3):
            sorted_vals = np.sort(points[:, dim])
            trimmed = sorted_vals[n_trim:-n_trim] if n_trim > 0 else sorted_vals
            center[dim] = np.mean(trimmed)
        return center
    
    source_center = trimmed_mean(source_pts, trim_percentage)
    pred_center = trimmed_mean(pred_pts, trim_percentage)
    translation = source_center - pred_center
    return pred_pts + translation, translation

def translate_median(source_pts, pred_pts):
    """
    Use median instead of mean for translation - more robust to outliers.
    """
    source_median = np.median(source_pts, axis=0)
    pred_median = np.median(pred_pts, axis=0)
    translation = source_median - pred_median
    return pred_pts + translation, translation

def translate_icp_based(source_pts, pred_pts, max_iterations=10, tolerance=1e-5):
    """
    Simplified ICP for translation only (no rotation).
    Iteratively finds correspondences and updates translation.
    """
    translated_pred = pred_pts.copy()
    translation = np.zeros(3)
    
    source_tree = cKDTree(source_pts)
    
    for iteration in range(max_iterations):
        # Find nearest neighbors
        distances, indices = source_tree.query(translated_pred)
        
        # Compute translation update
        matched_source = source_pts[indices]
        delta_translation = np.mean(matched_source - translated_pred, axis=0)
        
        # Update
        translated_pred += delta_translation
        translation += delta_translation
        
        # Check convergence
        if np.linalg.norm(delta_translation) < tolerance:
            break
    
    return translated_pred, translation

def main():
    args = get_args()
    tau = args.tau

    out_name = f"metric_plots_tau{tau}_translate{args.translate_method}_align{args.pose_align_method}_scale{args.scale_align_method}_ds{args.icp_downsample}_downsamplefactor{args.downsample_factor}"
    out_dir = os.path.join("viz", out_name)
    os.makedirs(out_dir, exist_ok=True)

    scenes = get_scenes_list(args)
    per_scene = []

    scale_dicts = {}
    if args.volume_method == 'log_all':
        dim_dict = {}
        dim_dict_path = os.path.join(out_dir, 'volume_dimensions.json')
        log_paths(args.log_file, ["vol dict", dim_dict_path])


    pdf_images = []

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pdf_path = os.path.join(args.pptx_path, f'{out_name}_{timestamp}.pdf')


    for idx, scene in tqdm(enumerate(scenes), desc="Processing scenes"):
        scene_dir = os.path.join(args.data_dir, 'scenes', scene)

        # --- load point clouds (your helpers) ---
        pcd_file = os.path.join(scene_dir, 'ns', 'point_cloud.ply')
        dt_file  = os.path.join(scene_dir, 'ns', 'dataparser_transforms.json')
        with open(os.path.join(scene_dir, 'features', f'voxel_size_{args.feature_load_name}.json'), 'r') as f:
            feature_voxel_size = json.load(f)['voxel_size']
        pts_a = load_ns_point_cloud(pcd_file, dt_file, ds_size=None)

        
        other_pcd_file = os.path.join(args.other_pc_dir, scene + '_00', f'ds_{args.downsample_factor}', 'point_cloud.ply')
        pts_b = load_ply_point_cloud(other_pcd_file)
        if args.translate_method == 'center':
            pts_a_translated = translate_to_origin(pts_a)
            pts_b_translated = translate_to_origin(pts_b)
            
        elif args.translate_method == 'nerf':
            pts_a_translated = pts_a
            # translate pts_b to pts_a centroid
            if idx not in STAY:
                pts_b_translated = pts_b - np.mean(pts_b, axis=0) + np.mean(pts_a, axis=0)
            else:
                pts_b_translated = pts_b
                
        elif args.translate_method == 'trimmed_mean_nerf':
            pts_a_translated = pts_a
            pts_b_translated, translation = translate_trimmed_mean(pts_a, pts_b, trim_percentage=0.5)
        elif args.translate_method == 'trimmed_median_nerf':
            pts_a_translated = pts_a
            pts_b_translated, translation = translate_median(pts_a, pts_b)
        elif args.translate_method == 'icp_nerf':
            pts_a_translated = pts_a
            pts_b_translated, translation = translate_icp_based(pts_a, pts_b, max_iterations=20, tolerance=1e-5)
        
        elif args.translate_method == 'nerf_self':
            pts_a_translated = pts_a
            pts_b_translated = pts_a
        else:
            pts_a_translated = pts_a
            pts_b_translated = pts_b
            #print(f"Scene {scene}: Applied trimmed mean translation {translation}")
        if args.pose_align_method != 'intensive':
            if args.pose_align_method == 'icp':
                #print('heehe')
                new_pts_b_list, [T], [stat] = icp_align_to_source(
                    source=pts_a_translated,
                    others=[pts_b_translated],
                    downsample_factor=args.icp_downsample
                )
                pts_b_rotated = new_pts_b_list[0]
            elif args.pose_align_method == 'manual_nerf':
                if idx in LIST_90:
                    angle = -90
                elif idx in LIST_60:
                    angle = -60
                else:
                    angle = 0
                
                R = rotation_matrix_z(angle)
                pts_b_rotated = (pts_b_translated @ R.T)

            # elif args.pose_align_method == 'real_world':
            #     info_dim = read_json(os.path.join(scene_dir, "info_dim.json"))
            #     dim_str = info_dim["dimension"]
            #     dim_llm = parse_dims_to_means(dim_str)
                
            #     # dim_gs_path = os.path.join(scene_dir, 'imsplat', 'sizes_gs.npy')
            #     # dim_gs = np.load(dim_gs_path)

            #     pose_dict = pca_pose(pts_b_translated)

            #     #centers = _to_numpy(centers)
            #     sizes, _, _ = size_from_pose_centers(pts_b_translated, pose=pose_dict)

            #     best = isotropic_scale(sizes, dim_llm)
            #     scale = best["s"]

            #     pts_b_scaled = pts_b_translated * scale

            # elif args.pose_align_method == 'nerf2physic_scale':
            #     diag_a = _aabb_diagonal_length(pts_a)
            #     diag_b = _aabb_diagonal_length(pts_b)

            #     eps = 1e-12
            #     if diag_b <= eps or diag_a <= eps:
            #         # Nothing sensible to scale; leave as-is but print a note so it's visible in logs.
            #         print(f"[nerf2physic_scale] Skipped scaling because diag_a={diag_a:.6g}, diag_b={diag_b:.6g}")
            #         pts_b_scaled = pts_b
            #         scale = 1.0
            #     else:
            #         scale = diag_a / (diag_b + eps)
            #         pts_b_scaled = pts_b * scale
            #         print(f"[nerf2physic_scale] Applied isotropic scale={scale:.6g} (diag_a={diag_a:.6g}, diag_b={diag_b:.6g})")

            #     scale_dicts[scene] = scale
            # elif args.pose_align_method == 'manual':
            #     if idx in LIST_90:
            #         angle = -90
            #     elif idx in LIST_60:
            #         angle = -60
            #     else:
            #         angle = 0
                
            #     R = rotation_matrix_y(angle)
            #     pts_b_scaled = (pts_b_translated @ R.T)

            # else:
            #     pts_b_scaled = pts_b_translated
        else:
            pts_b_rotated, T_total, T_pca, T_icp = intensive_alignment(
                source=pts_a_translated,
                other=pts_b_translated,
                downsample_factor=args.icp_downsample,
                do_icp=True,
                icp_trim=0.2,
                icp_max_iters=40
            )
        
        if args.scale_align_method == 'nerf_scale':
            diag_a = _aabb_diagonal_length(pts_a_translated)
            diag_b = _aabb_diagonal_length(pts_b_rotated)

            eps = 1e-12
            if diag_b <= eps or diag_a <= eps:
                # Nothing sensible to scale; leave as-is but print a note so it's visible in logs.
                print(f"[nerf2physic_scale] Skipped scaling because diag_a={diag_a:.6g}, diag_b={diag_b:.6g}")
                scale = 1.0
            else:
                scale = diag_a / (diag_b + eps)
                pts_b_scaled = pts_b_rotated * scale
                print(f"[nerf2physic_scale] Applied isotropic scale={scale:.6g} (diag_a={diag_a:.6g}, diag_b={diag_b:.6g})")

            scale_dicts[scene] = scale
        elif args.scale_align_method == 'real_world':
            info_dim = read_json(os.path.join(scene_dir, "info_dim.json"))
            dim_str = info_dim["dimension"]
            dim_llm = parse_dims_to_means(dim_str)
            
            # dim_gs_path = os.path.join(scene_dir, 'imsplat', 'sizes_gs.npy')
            # dim_gs = np.load(dim_gs_path)

            pose_dict = pca_pose(pts_b_rotated)

            #centers = _to_numpy(centers)
            sizes, _, _ = size_from_pose_centers(pts_b_rotated, pose=pose_dict)

            best = isotropic_scale(sizes, dim_llm)
            scale = best["s"]

            pts_b_scaled = pts_b_rotated * scale

        elif args.scale_align_method == 'nerf2physic_scale':
            diag_a = _aabb_diagonal_length(pts_a)
            diag_b = _aabb_diagonal_length(pts_b)

            eps = 1e-12
            if diag_b <= eps or diag_a <= eps:
                # Nothing sensible to scale; leave as-is but print a note so it's visible in logs.
                print(f"[nerf2physic_scale] Skipped scaling because diag_a={diag_a:.6g}, diag_b={diag_b:.6g}")
                pts_b_scaled = pts_b_rotated
                scale = 1.0
            else:
                scale = diag_a / (diag_b + eps)
                pts_b_scaled = pts_b_rotated * scale
                print(f"[nerf2physic_scale] Applied isotropic scale={scale:.6g} (diag_a={diag_a:.6g}, diag_b={diag_b:.6g})")

            scale_dicts[scene] = scale


        else:
            pts_b_scaled = pts_b_rotated

        # translated then rotated then scaled

        # --- load normals if you have them for BOTH sets (optional for NC) ---
        normals_a = _safe_normals_load(os.path.join(scene_dir, 'ns', 'point_cloud_normals.txt'))
        normals_b = _safe_normals_load(os.path.join(args.other_pc_dir, scene + '_00', 'point_cloud_normals.txt'))

        m = compute_scene_metrics(pts_a_translated, pts_b_scaled, normals_a, normals_b, tau=tau)
        m["scene"] = scene
        bad_fscore = False
        if m['fscore'] <= args.f_threshold and args.f_threshold > 0:
            bad_fscore = True
        per_scene.append(m)

        if args.volume_method == 'log_all':
            # nerf dim
            pose_dict_nerf = pca_pose(pts_a)

                #centers = _to_numpy(centers)
            dim_nerf, _, _ = size_from_pose_centers(pts_a, pose=pose_dict_nerf)

            # physx dim
            pose_dict_physx = pca_pose(pts_b)
            dim_physx, _, _ = size_from_pose_centers(pts_b, pose=pose_dict_physx)

            dim_dict[scene] = {
                'dim_nerf': dim_nerf.tolist(),
                'dim_physx': dim_physx.tolist(),
            }
        if args.plot_type == 'pc':
            # plot 2 sets of point cloud, and concatenate them side by side
            save_point_cloud_plot(pts_a, os.path.join('viz', scene, out_name, 'nerf_pc.png'), point_size=3.0, axes_mode='fixed')
            save_point_cloud_plot(pts_b, os.path.join('viz', scene, out_name, 'physx_pc.png'), point_size=3.0, axes_mode='fixed')
            save_point_cloud_plot(pts_b_scaled, os.path.join('viz', scene, out_name, 'physx_pc_scaled.png'), point_size=3.0, axes_mode='fixed')

            orig_rgb_image_path = os.path.join(scene_dir, 'images', f'{scene}_00.png')
            #print('orig_rgb_image_path', orig_rgb_image_path)
            combine_save_path = os.path.join('viz', scene, out_name, 'comparison.png')
            # stack_images(
            #     images_paths_list=[
            #         orig_rgb_image_path,
            #         os.path.join('viz', scene, out_name,'nerf_pc.png'),
            #         os.path.join('viz', scene, out_name,'physx_pc.png')
            #     ],
            #     save_path=combine_save_path
            # )

            combined_img = stack_images_pil(
                candidate_paths=[orig_rgb_image_path,
                                 os.path.join('viz', scene, out_name,'nerf_pc.png'),
                                os.path.join('viz', scene, out_name,'physx_pc.png'),
                                os.path.join('viz', scene, out_name,'physx_pc_scaled.png')],
                save_path=combine_save_path
            )
            #if bad_fscore:
            log_paths(args.log_file, [combine_save_path])

            if bad_fscore:
                log_paths(args.log_file, ['bad fscore', combine_save_path])

            pdf_images.append(combined_img)

            # blank_slide_layout = prs.slide_layouts[6]

            # slide = prs.slides.add_slide(blank_slide_layout)
            # slide.shapes.add_picture(combine_save_path, Inches(0), Inches(0), width=Inches(10))  # Adjust size if needed
    
        if args.save_pc:
            save_path = os.path.join(args.data_dir, 'scenes', scene, 'physx', f'point_cloud_{args.translate_method}.ply')
            os.makedirs(os.path.dirname(save_path), exist_ok=True)
            if args.verbose:
                print('num pts saved:', pts_b_scaled.shape)
            save_ply_point_cloud(pts_b_scaled, save_path)

    if args.scale_align_method == 'nerf2physic_scale' and args.overwrite:
        write_json(scale_dicts, args.nerf_vs_physx_scale)

    if args.volume_method == 'log_all':
        write_json(dim_dict, dim_dict_path)
    # --- save CSV ---
    csv_path = os.path.join(out_dir, "metrics_per_scene.csv")
    save_metrics_csv(per_scene, csv_path)

    mean_chamfer = np.mean([m["chamfer"] for m in per_scene])
    mean_fscore  = np.mean([m["fscore"]  for m in per_scene])
    
    columns = {
        'dataset': 'abo500_test',
        'tau': tau,
        'translate_method': args.translate_method,
        'pose_align_method': args.pose_align_method,
        'scale_align_method': args.scale_align_method,
        'chamfer': mean_chamfer,
        'fscore': mean_fscore
    }
    
    log_eval_result(columns, csv_path=args.result_path)
    # --- plots (one figure per metric) ---
    labels = [str(i) for i in range(len(per_scene))]
    chamfers = [m["chamfer"] for m in per_scene]
    fscores  = [m["fscore"]  for m in per_scene]
    ncs      = [m["normal_consistency"] if m["normal_consistency"] is not None else np.nan for m in per_scene]

    plot_series(labels, chamfers, "Chamfer Distance per Scene", "Chamfer (↓)", os.path.join(out_dir, "chamfer_per_scene.png"))
    plot_series(labels, fscores,  f"F-score@τ per Scene (τ={tau})", "F-score (↑)", os.path.join(out_dir, "fscore_per_scene.png"))
    plot_series(labels, ncs,      f"Normal Consistency per Scene (τ={tau})", "Normal Consistency (↑)", os.path.join(out_dir, "normal_consistency_per_scene.png"))

    print(f"Saved CSV: {csv_path}")
    print(f"Saved plots to: {out_dir}")

    pdf_images[0].save(pdf_path, save_all=True, append_images=pdf_images[1:], resolution=100.0, quality=95)
    print(f"Saved PDF: {pdf_path}")
    log_paths(args.log_file, ["pdf", pdf_path])

if __name__ == "__main__":
    main()
